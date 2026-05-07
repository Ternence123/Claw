#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""在SPC培训PPT第3页位置插入名词解释页（含σ̂ S̄ R̄ X̄̄ Σ）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shutil
import zipfile
import lxml.etree as ET_lxml
import os

INPUT  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUTPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_名词解释.pptx"

GREEN = RGBColor(0x2E, 0x8B, 0x57)
BLUE  = RGBColor(0x1E, 0x3A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREEN= RGBColor(0xE8, 0xF5, 0xE9)
LBLUE = RGBColor(0xE3, 0xF2, 0xFD)
LORANGE=RGBColor(0xFF, 0xF8, 0xE1)
DARK  = RGBColor(0x21, 0x21, 0x21)


def clear_placeholders(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.text = ""
        if shape.is_placeholder:
            shape.element.getparent().remove(shape.element)


def make_term_slide(prs):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    clear_placeholders(slide)

    # ── 顶部绿色标题条 ──────────────────────────
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = GREEN
    rect.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "📖  SPC 核心名词与符号解释"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ── 术语列表（17项，按逻辑分组）────────────────────
    terms = [
        # (符号, 名称, 含义, 公式)
        ("Σ (Sigma/求和)", "求和符号  Summation",
         "对一系列数值进行累加运算的符号",
         "Σxᵢ = x₁ + x₂ + … + xₙ"),

        ("x̄ (X-Bar)", "样本均值  Sample Mean",
         "单个子组内数据的平均值，反映该子组中心",
         "x̄ = Σxᵢ / n"),

        ("X̄̄ (X-Double-Bar)", "总均值  Grand Mean",
         "所有子组均值 x̄ 的总平均，即过程中心线 CL",
         "X̄̄ = (x̄₁ + x̄₂ + … + x̄k) / k"),

        ("R", "极差  Range",
         "单个子组内最大值与最小值之差",
         "R = xₘₐₓ − xₘᵢₙ"),

        ("R̄ (R-Bar)", "平均极差  Average Range",
         "所有子组极差 R 的平均值，用于估算过程变异",
         "R̄ = (R₁ + R₂ + … + Rk) / k"),

        ("σ (Sigma)", "总体标准差  Population StdDev",
         "描述总体数据离散程度的参数（常用希腊字母σ）",
         "σ = √[ Σ(xᵢ−μ)² / N ]"),

        ("σ̂ (Sigma Hat)", "估计标准差  Estimated StdDev",
         "用样本数据估算的总体标准差（σ的估计值）",
         "σ̂ = R̄ / d₂   （用极差估算）"),

        ("s (或 S)", "样本标准差  Sample StdDev",
         "单个子组内用样本公式算出的标准差",
         "s = √[ Σ(xᵢ−x̄)² / (n−1) ]"),

        ("S̄ (S-Bar)", "平均样本标准差",
         "所有子组样本标准差 s 的平均值",
         "S̄ = (s₁ + s₂ + … + sₖ) / k"),

        ("USL", "规格上限  Upper Spec Limit",
         "客户/设计允许的最大值，超出即不合格",
         "由产品设计决定（非统计值）"),

        ("LSL", "规格下限  Lower Spec Limit",
         "客户/设计允许的最小值，低于即不合格",
         "由产品设计决定（非统计值）"),

        ("UCL", "上控制限  Upper Control Limit",
         "控制图预警上限；超出说明过程异常",
         "UCL = X̄̄ + A₂·R̄"),

        ("CL", "中心线  Center Line",
         "控制图基准线，取总均值 X̄̄",
         "CL = X̄̄"),

        ("LCL", "下控制限  Lower Control Limit",
         "控制图预警下限；低于说明过程异常",
         "LCL = X̄̄ − A₂·R̄"),

        ("Cp", "过程能力指数  Process Capability",
         "衡量过程潜在能力（假设均值无偏移）",
         "Cp = (USL − LSL) / (6σ̂)"),

        ("Cpk", "实际过程能力指数",
         "衡量过程实际能力（考虑均值偏移）",
         "Cpk = min[(USL−μ)/3σ̂ , (μ−LSL)/3σ̂]"),

        ("3σ 原理", "三西格玛法则  Three-Sigma Rule",
         "正态分布中 99.73% 的数据落在均值 ±3σ 内",
         "P(μ ± 3σ) = 99.73%"),
    ]

    # ── 三列卡片布局 ─────────────────────────────
    col_x_list = [Inches(0.2), Inches(4.45), Inches(8.7)]
    card_w     = Inches(4.15)
    start_y    = Inches(1.25)
    row_h      = Inches(0.88)

    for i, (sym, name, meaning, formula) in enumerate(terms):
        col_idx = i // 6          # 每列6个
        row_idx = i % 6
        if col_idx >= 3:
            break   # 最多3列×6行=18个
        cx = col_x_list[col_idx]
        cy = start_y + row_idx * row_h

        # 卡片背景（三色交替：绿/蓝/橙）
        card = slide.shapes.add_shape(1, cx, cy, card_w, row_h - Inches(0.05))
        card.fill.solid()
        if col_idx == 0:
            card.fill.fore_color.rgb = LGREEN
        elif col_idx == 1:
            card.fill.fore_color.rgb = LBLUE
        else:
            card.fill.fore_color.rgb = LORANGE
        card.line.color.rgb = GREEN
        card.line.width = Pt(0.5)

        # 符号（加粗）
        sym_box = slide.shapes.add_textbox(
            cx + Inches(0.1), cy + Inches(0.04),
            Inches(1.5), Inches(0.26))
        p = sym_box.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN

        # 名称
        name_box = slide.shapes.add_textbox(
            cx + Inches(1.65), cy + Inches(0.04),
            Inches(2.4), Inches(0.26))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(8.5)
        p.font.color.rgb = DARK

        # 含义
        mean_box = slide.shapes.add_textbox(
            cx + Inches(0.1), cy + Inches(0.32),
            card_w - Inches(0.15), Inches(0.26))
        p = mean_box.text_frame.paragraphs[0]
        p.text = "📌 " + meaning
        p.font.size = Pt(8)
        p.font.color.rgb = DARK

        # 公式
        if formula:
            fbox = slide.shapes.add_textbox(
                cx + Inches(0.1), cy + Inches(0.58),
                card_w - Inches(0.15), Inches(0.25))
            p = fbox.text_frame.paragraphs[0]
            p.text = "🔢 " + formula
            p.font.size = Pt(7.5)
            p.font.italic = True
            p.font.color.rgb = BLUE

    return slide


def reorder_slides_zip(input_path, output_path, insert_pos=2):
    """将最后一页移到 insert_pos 位置"""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    shutil.copy(input_path, output_path)

    with zipfile.ZipFile(output_path, 'r') as z:
        pres_bytes = z.read('ppt/presentation.xml')

    root = ET_lxml.fromstring(pres_bytes)
    sld_id_lst = root.find('p:sldIdLst', ns)
    sld_ids = sld_id_lst.findall('p:sldId', ns)

    total = len(sld_ids)
    if insert_pos >= total:
        print(f"⚠  insert_pos={insert_pos} >= 总页数{total}，不调整顺序")
        return

    last_elem = sld_ids[-1]
    sld_id_lst.remove(last_elem)
    sld_id_lst.insert(insert_pos, last_elem)

    new_pres_bytes = ET_lxml.tostring(root, xml_declaration=True, encoding='UTF-8')

    tmp_zip = output_path + '.tmp.zip'
    with zipfile.ZipFile(output_path, 'r') as zin:
        with zipfile.ZipFile(tmp_zip, 'w') as zout:
            for item in zin.namelist():
                if item == 'ppt/presentation.xml':
                    zout.writestr(item, new_pres_bytes)
                else:
                    zout.writestr(item, zin.read(item))
    os.replace(tmp_zip, output_path)
    print(f"✅ 页码调整完成：名词解释页已移至第{insert_pos+1}页位置")


def main():
    prs = Presentation(INPUT)
    orig_count = len(prs.slides)
    print(f"原文件页数: {orig_count}")

    make_term_slide(prs)

    tmp_file = INPUT + '.tmp.pptx'
    prs.save(tmp_file)
    print(f"已追加名词解释页（共17个符号），正在调整页码顺序...")

    reorder_slides_zip(tmp_file, OUTPUT, insert_pos=2)
    os.remove(tmp_file)

    # 验证
    prs2 = Presentation(OUTPUT)
    print(f"\n✅ 已生成: {OUTPUT}")
    print(f"   总页数: {len(prs2.slides)}")
    for i, s in enumerate(prs2.slides, 1):
        title = ''
        for sh in s.shapes:
            if sh.has_text_frame and sh.text.strip():
                title = sh.text.strip()[:40]
                break
        mark = "  ← 📖 名词解释" if ("名词解释" in title or "SPC 核心" in title) else ""
        print(f"   第{i}页: {title}{mark}")

if __name__ == "__main__":
    main()
