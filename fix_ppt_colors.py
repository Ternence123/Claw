"""
SPC PPT 配色统一 + 布局平铺修复脚本 v2
方案：绿色(#10B981) + 蓝色(#0EA5E9) 为主色调
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import copy, os

PPT_PATH = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUT_PATH  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v3.pptx"

# ── 新配色 ─────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x29, 0x3B)   # 深蓝灰 标题背景
BLUE_MAIN  = RGBColor(0x0E, 0xA5, 0xE9)   # 主蓝 色块
GREEN_MAIN = RGBColor(0x10, 0xB9, 0x81)   # 主绿 色块
BLUE_LIGHT = RGBColor(0xF0, 0xF9, 0xFF)   # 浅蓝 浅背景
GREEN_LIGHT= RGBColor(0xEC, 0xFD, 0xF5)   # 浅绿 浅背景
RED_WARN   = RGBColor(0xEF, 0x44, 0x44)   # 红色 警告/重要
YELLOW_HL  = RGBColor(0xF5, 0x9E, 0x0B)   # 黄色 高亮
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)   # 深灰文字（浅底）
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)

def hex_from_rgb(rgb):
    """RGBColor → hex字符串（大写，6位）"""
    try:
        s = str(rgb)
        if len(s) == 6:
            return s.upper()
        # 有些版本返回 0xXXXXXX
        return s.replace('0x','').upper().zfill(6)
    except:
        return ""

def is_dark(hex6):
    """YIQ法判断颜色深浅，hex6如 '1E293B'"""
    try:
        r = int(hex6[0:2], 16)
        g = int(hex6[2:4], 16)
        b = int(hex6[4:6], 16)
        yiq = (r * 299 + g * 587 + b * 114) / 1000
        return yiq < 128
    except:
        return True

def apply_fill(shape, color):
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    except:
        pass

def apply_text_rgb(text_frame, color_rgb):
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = color_rgb
            except:
                pass

def fix_shape(shape):
    """递归处理形状：改填充色 + 改文字颜色"""
    # 1) 形状填充
    fill_hex = None
    if hasattr(shape, 'fill') and shape.has_text_frame is not None:
        try:
            if shape.fill.type is not None:
                rgb = shape.fill.fore_color.rgb
                fill_hex = hex_from_rgb(rgb)
        except:
            pass

    # 颜色映射
    new_color = None
    if fill_hex:
        m = {
            # 蓝色系 → 主蓝
            "0891B2": BLUE_MAIN, "2563EB": BLUE_MAIN, "2E86AB": BLUE_MAIN,
            # 深蓝/黑 → 深蓝背景
            "1E3A5F": BLUE_DARK, "1E293B": BLUE_DARK, "0D2137": BLUE_DARK,
            # 绿色系 → 主绿
            "27AE60": GREEN_MAIN, "10B981": GREEN_MAIN,
            # 红色 → 警告红
            "EF4444": RED_WARN, "C0392B": RED_WARN,
            # 黄色 → 高亮黄
            "F59E0B": YELLOW_HL, "E07B39": YELLOW_HL,
            # 紫色 → 主蓝（去掉紫）
            "8E44AD": BLUE_MAIN, "7C3AED": BLUE_MAIN,
            # 浅灰/白背景
            "F1F5F9": BLUE_LIGHT, "FFFFF0": BLUE_LIGHT, "FFFFFF": WHITE,
            "CBD5E1": BLUE_LIGHT,
        }
        new_color = m.get(fill_hex)

    if new_color:
        apply_fill(shape, new_color)
        # 2) 文字颜色：深色背景 → 白色；浅色背景 → 深灰
        if shape.has_text_frame:
            text_color = WHITE if is_dark(hex_from_rgb(new_color)) else DARK_TEXT
            apply_text_rgb(shape.text_frame, text_color)
    else:
        # 无填充的形状：若含文字，确保文字颜色合适
        if shape.has_text_frame:
            apply_text_rgb(shape.text_frame, DARK_TEXT)

    # 递归处理组合形状
    try:
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                fix_shape(sub)
    except:
        pass

def fix_slide(slide):
    for shape in slide.shapes:
        fix_shape(shape)

def main():
    prs = Presentation(PPT_PATH)
    print(f"共 {len(prs.slides)} 张幻灯片，开始修复配色...")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  幻灯片 {i}...")
        fix_slide(slide)
    prs.save(OUT_PATH)
    print(f"\n✅ 已保存至：{OUT_PATH}")

if __name__ == "__main__":
    main()
