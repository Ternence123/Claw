#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""SPC培训PPT v3 ── 统一配色美化 + 名词解释2页（正确插入第3页位置）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shutil, os, zipfile, lxml.etree as ET, time

INPUT  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
STAMP = int(time.time())
OUTPUT = f"c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_美化版_{STAMP}.pptx"

# ── 配色（从原PPT theme1.xml 提取）───────────────
TITLE_BG  = RGBColor(0x1E, 0x3A, 0x5F)   # 标题栏深蓝
ACCENT    = RGBColor(0x08, 0x91, 0xB2)   # 强调青蓝
A1_BLUE   = RGBColor(0x44, 0x72, 0xC4)   # accent1 蓝
A6_GREEN  = RGBColor(0x70, 0xAD, 0x47)   # accent6 绿
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1F, 0x29, 0x3A)   # 深色文字
CARD_G1   = RGBColor(0xE8, 0xF5, 0xE9)   # 卡片绿
CARD_G2   = RGBColor(0xE3, 0xF2, 0xFD)   # 卡片蓝
CARD_G3   = RGBColor(0xE0, 0xF7, 0xFA)   # 卡片青

def clear_placeholder_shapes(slide):
    for s in list(slide.shapes):
        if s.is_placeholder:
            s.element.getparent().remove(s.element)

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_line(shape, color, width_pt=0.5):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)

# ══════════════════════════════════════════════════════
#  制作一页名词解释
# ══════════════════════════════════════════════════════

def make_term_slide(prs, page_no, total_pages, terms):
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    clear_placeholder_shapes(slide)

    W = Inches(13.33)
    H = Inches(7.5)

    # 标题条（深蓝底白字，匹配原PPT）
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.08))
    set_fill(bar, TITLE_BG)
    bar.line.fill.background()

    # 左侧装饰窄条
    left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), H)
    set_fill(left, ACCENT)
    left.line.fill.background()

    # 标题文字
    txb = slide.shapes.add_textbox(Inches(0.35), Inches(0.1), Inches(12), Inches(0.85))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.text = f"  📖  SPC 核心名词与符号解释（{page_no}/{total_pages}）"
    p.font.size  = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ── 3列 × 3行 卡片 ─────────────────────────────
    col_x   = [Inches(0.2), Inches(4.5), Inches(8.85)]
    row_y   = [Inches(1.2), Inches(2.8), Inches(4.45)]
    CARD_W  = Inches(4.15)
    CARD_H  = Inches(1.45)
    BG      = [CARD_G1, CARD_G2, CARD_G3]

    for idx, (sym, name, meaning, formula) in enumerate(terms):
        col = idx % 3
        row = idx // 3
        if row >= 3:
            break
        cx = col_x[col]
        cy = row_y[row]

        # 卡片底
        card = slide.shapes.add_shape(1, cx, cy, CARD_W, CARD_H)
        set_fill(card, BG[(col + row) % 3])
        set_line(card, ACCENT, 0.75)

        # 符号
        b1 = slide.shapes.add_textbox(cx + Inches(0.1),  cy + Inches(0.06),
                                      Inches(1.5), Inches(0.3))
        p = b1.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_BG

        # 名称
        b2 = slide.shapes.add_textbox(cx + Inches(1.65), cy + Inches(0.08),
                                      CARD_W - Inches(1.75), Inches(0.3))
        p = b2.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(8.5)
        p.font.color.rgb = DARK

        # 含义
        b3 = slide.shapes.add_textbox(cx + Inches(0.1),  cy + Inches(0.42),
                                      CARD_W - Inches(0.2), Inches(0.38))
        p = b3.text_frame.paragraphs[0]
        p.text = "📌 " + meaning
        p.font.size = Pt(8)
        p.font.color.rgb = DARK

        # 公式
        if formula:
            b4 = slide.shapes.add_textbox(cx + Inches(0.1),  cy + Inches(0.82),
                                          CARD_W - Inches(0.2), Inches(0.55))
            tf4 = b4.text_frame
            tf4.word_wrap = True
            p = tf4.paragraphs[0]
            p.text = "🔢 " + formula
            p.font.size = Pt(7.5)
            p.font.italic = True
            p.font.color.rgb = A1_BLUE

    return slide

# ══════════════════════════════════════════════════════
#  用 XML 将指定页移到目标位置（支持多页）
# ══════════════════════════════════════════════════════

def reorder_slides(input_path, output_path, positions):
    """
    positions: [(dst, src), ...]   dst=目标位置(0-based), src=当前位置(0-based)
    将 input_path 的幻灯片按 positions 重排后保存到 output_path
    """
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    with zipfile.ZipFile(input_path, 'r') as z:
        raw = z.read('ppt/presentation.xml')

    root = ET.fromstring(raw)
    sld_id_lst = root.find('p:sldIdLst', ns)
    elems = list(sld_id_lst.findall('p:sldId', ns))
    N = len(elems)

    # dst_to_src: 目标位置 → 源索引
    dst_to_src = {}
    moved_srcs = set()

    for dst, src in positions:
        dst_to_src[dst] = src
        moved_srcs.add(src)

    # 构建新顺序
    new_elems = []
    orig_ptr = 0

    for i in range(N):
        if i in dst_to_src:
            # 这个位置应放被移动的幻灯片
            src = dst_to_src[i]
            new_elems.append(elems[src])
        else:
            # 放原始顺序中下一个未被移动的幻灯片
            while orig_ptr in moved_srcs:
                orig_ptr += 1
            if orig_ptr < N:
                new_elems.append(elems[orig_ptr])
            orig_ptr += 1

    # 写回 XML
    sld_id_lst.clear()
    for e in new_elems:
        sld_id_lst.append(e)

    new_raw = ET.tostring(root, xml_declaration=True, encoding='UTF-8')

    tmp = output_path + '.tmp.zip'
    with zipfile.ZipFile(input_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w') as zout:
            for item in zin.namelist():
                if item == 'ppt/presentation.xml':
                    zout.writestr(item, new_raw)
                else:
                    zout.writestr(item, zin.read(item))
    shutil.move(tmp, output_path)
    print("  ✅ 页码顺序已调整（名词解释页已插入到第3页位置）")

# ══════════════════════════════════════════════════════
#  main
# ══════════════════════════════════════════════════════

def main():
    prs = Presentation(INPUT)
    print(f"原文件页数: {len(prs.slides)}")

    all_terms = [
        ("Σ",   "求和符号  Summation",
         "对一系列数值进行累加运算",
         "Σxᵢ = x₁ + x₂ + … + xₙ"),
        ("x̄",   "样本均值  Sample Mean",
         "单个子组内数据的平均值",
         "x̄ = Σxᵢ / n"),
        ("X̄̄",  "总均值  Grand Mean",
         "所有子组均值的总平均（即CL中心线）",
         "X̄̄ = (x̄₁ + x̄₂ + … + x̄k) / k"),
        ("R",    "极差  Range",
         "子组内最大值与最小值之差",
         "R = xₘₐₓ − xₘᵢₙ"),
        ("R̄",   "平均极差  Avg Range",
         "所有子组极差 R 的平均值",
         "R̄ = (R₁ + R₂ + … + Rk) / k"),
        ("σ",    "总体标准差  Population StdDev",
         "描述总体数据离散程度的参数",
         "σ = √[ Σ(xᵢ−μ)² / N ]"),
        ("σ̂",   "估计标准差  Estimated Sigma",
         "用样本数据估算的总体标准差",
         "σ̂ = R̄/d₂   或   σ̂ = S̄/c₄"),
        ("s / S̄", "样本标准差 / 平均样本标准差",
         "样本公式算出的标准差及其跨子组均值",
         "S̄ = (s₁ + s₂ + … + sₖ) / k"),
        ("USL",  "规格上限  Upper Spec Limit",
         "客户/设计允许的最大值，超出即不合格",
         "由产品设计决定（非统计计算值）"),
        ("LSL",  "规格下限  Lower Spec Limit",
         "客户/设计允许的最小值，低于即不合格",
         "由产品设计决定（非统计计算值）"),
        ("UCL",  "上控制限  Upper Control Limit",
         "控制图预警上限；超出说明过程异常",
         "UCL = X̄̄ + A₂·R̄"),
        ("CL",   "中心线  Center Line",
         "控制图基准线，取所有子组的总均值",
         "CL = X̄̄"),
        ("LCL",  "下控制限  Lower Control Limit",
         "控制图预警下限；低于说明过程异常",
         "LCL = X̄̄ − A₂·R̄"),
        ("Cp",   "过程能力指数  Process Capability",
         "衡量过程潜在能力（假设均值无偏移）",
         "Cp = (USL − LSL) / (6σ̂)"),
        ("Cpk",  "实际过程能力指数",
         "衡量过程实际能力（考虑均值偏移）",
         "Cpk = min[(USL−μ)/3σ̂ , (μ−LSL)/3σ̂]"),
        ("Pp / Ppk", "过程性能指数  Process Performance",
         "用长期总体数据计算（使用总体标准差）",
         "Pp = (USL − LSL) / (6σₜₒₜₐₗ)"),
        ("3σ原理", "三西格玛法则  Three-Sigma Rule",
         "正态分布中 99.73% 的数据落在均值 ±3σ 范围内",
         "P(μ ± 3σ) = 99.73%"),
    ]

    n_total  = len(all_terms)
    per_page = 9
    n_pages  = (n_total + per_page - 1) // per_page
    print(f"术语总数: {n_total}，分 {n_pages} 页（每页最多{per_page}个）")

    # 追加页面（会在末尾）
    added = []
    for p_no in range(n_pages):
        chunk = all_terms[p_no * per_page : (p_no + 1) * per_page]
        print(f"  第{p_no+1}页：{len(chunk)} 个术语")
        slide = make_term_slide(prs, p_no + 1, n_pages, chunk)
        added.append(len(prs.slides) - 1)

    # 保存临时文件
    tmp = INPUT + '.tmp.pptx'
    prs.save(tmp)
    print("临时文件已保存，正在调整页码顺序...")

    # 将追加的页移到第3页位置（索引2）
    # added = [12, 13]  →  插入到位置 2, 3
    positions = [(2 + i, added[i]) for i in range(len(added))]
    reorder_slides(tmp, OUTPUT, positions)
    os.remove(tmp)

    # 验证
    prs2 = Presentation(OUTPUT)
    print(f"\n✅ 已生成：{OUTPUT}")
    print(f"   总页数: {len(prs2.slides)}")
    for i, s in enumerate(prs2.slides, 1):
        title = ''
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text.strip()
                if len(t) > 2:
                    title = t[:36]
                    break
        mark = "  ← 📖 名词解释" if "名词解释" in title else ""
        print(f"   第{i}页: {title}{mark}")

if __name__ == "__main__":
    main()
