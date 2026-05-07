"""
PPT 配色修复 v4
1. 红/黄色块 → 尽量改为蓝/绿（只保留很小的装饰色块）
2. 文字颜色：若文字框位于深色背景上 → 白色，否则深灰色
3. 处理"独立文本框位于色块上方"的情况
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

PPT_PATH = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUT_PATH  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v3.pptx"

# ── 目标配色 ─────────────────────────────────
C_DARK_BG    = RGBColor(0x1E, 0x29, 0x3B)   # 深蓝灰 标题栏
C_BLUE_MAIN  = RGBColor(0x0E, 0xA5, 0xE9)   # 主蓝
C_GREEN_MAIN = RGBColor(0x10, 0xB9, 0x81)   # 主绿
C_BLUE_LIGHT = RGBColor(0xF0, 0xF9, 0xFF)   # 浅蓝背景
C_GREEN_LIGHT= RGBColor(0xEC, 0xFD, 0xF5)   # 浅绿背景
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)   # 深灰文字
C_MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)
C_RED_WARN   = RGBColor(0xEF, 0x44, 0x44)   # 保留少量红色
C_YELLOW_HL = RGBColor(0xF5, 0x9E, 0x0B)   # 保留少量黄色

def hex6(rgb_color):
    try:
        s = str(rgb_color).upper().replace("0X", "")
        return s.zfill(6)
    except:
        return ""

def is_dark(hex6_str):
    try:
        r=int(hex6_str[0:2],16); g=int(hex6_str[2:4],16); b=int(hex6_str[4:6],16)
        return (r*299+g*587+b*114)/1000 < 128
    except:
        return True

def get_fill_hex(shape):
    """获取形状填充色 hex6，无填充返回 None"""
    try:
        if shape.fill and shape.fill.type is not None:
            return hex6(shape.fill.fore_color.rgb)
    except:
        pass
    return None

def set_fill(shape, color):
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    except:
        pass

def set_text_rgb(text_frame, color_rgb):
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = color_rgb
            except:
                pass

def remap(old_hex):
    """
    颜色重映射规则：
    - 蓝系 → 主蓝
    - 深蓝/黑 → 深蓝背景
    - 绿系 → 主绿
    - 红/紫系 → 主蓝（大幅减少红紫使用）
    - 黄系 → 主绿（大幅减少黄使用）
    - 浅灰/白 → 浅蓝背景或白色
    """
    m = {
        "0891B2":"BLUE", "2563EB":"BLUE", "2E86AB":"BLUE",
        "1E3A5F":"DARK_BG", "1E293B":"DARK_BG", "0D2137":"DARK_BG",
        "27AE60":"GREEN", "10B981":"GREEN",
        # 红 → 主蓝（去掉大部分红）
        "EF4444":"BLUE", "C0392B":"BLUE",
        # 黄 → 主绿（去掉大部分黄）
        "F59E0B":"GREEN", "E07B39":"GREEN",
        # 紫 → 主蓝
        "8E44AD":"BLUE", "7C3AED":"BLUE",
        # 浅灰
        "F1F5F9":"BLUE_LIGHT", "CBD5E1":"BLUE_LIGHT",
        "FFFFF0":"BLUE_LIGHT", "FFFFFF":"WHITE",
    }
    cat = m.get(old_hex)
    if cat == "BLUE":     return C_BLUE_MAIN
    if cat == "GREEN":   return C_GREEN_MAIN
    if cat == "DARK_BG": return C_DARK_BG
    if cat == "BLUE_LIGHT": return C_BLUE_LIGHT
    if cat == "WHITE":   return C_WHITE
    return None

def get_shape_rect(shape):
    """返回 (left, top, right, bottom) 英寸"""
    try:
        l = shape.left.inches
        t = shape.top.inches
        return (l, t, l + shape.width.inches, t + shape.height.inches)
    except:
        return None

def rect_contains(outer, inner):
    """判断 outer 矩形是否包含 inner 矩形（允许少量边距）"""
    if not outer or not inner:
        return False
    ol, ot, or_, ob = outer
    il, it, ir, ib = inner
    # inner 完全在 outer 内部（允许 0.05 英寸误差）
    return (il >= ol - 0.05) and (it >= ot - 0.05) and (ir <= or_ + 0.05) and (ib <= ob + 0.05)

def process_slide(slide):
    """处理一张幻灯片：先改填充色，再修正文字颜色"""
    
    # 1) 收集所有色块（有填充色的形状）
    color_blocks = []  # [(shape, fill_hex, rect)]
    for shape in slide.shapes:
        fh = get_fill_hex(shape)
        if fh:
            rect = get_shape_rect(shape)
            color_blocks.append((shape, fh, rect))
        # 递归处理组合形状
        try:
            if shape.shape_type == 6:  # GROUP
                for sub in shape.shapes:
                    fh = get_fill_hex(sub)
                    if fh:
                        rect = get_shape_rect(sub)
                        color_blocks.append((sub, fh, rect))
        except:
            pass
    
    # 2) 修改所有色块的填充色
    for item in color_blocks:
        shape, old_hex, rect = item
        new_color = remap(old_hex)
        if new_color:
            set_fill(shape, new_color)
            # 更新记录中的颜色（用于后续文字颜色判断）
            new_hex = hex6(new_color)
            item = (shape, new_hex, rect)
    
    # 3) 处理文字颜色：
    #    - 对于自身有填充的形状：根据填充色深浅设文字色
    #    - 对于无填充的文字框：找到它下方的色块，根据色块颜色设文字色
    
    # 收集所有文字框（自身无填充，或有填充但也要处理）
    text_items = []  # [(shape, rect)]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            rect = get_shape_rect(shape)
            text_items.append((shape, rect))
            # 递归处理组合
            try:
                if shape.shape_type == 6:
                    for sub in shape.shapes:
                        if sub.has_text_frame and sub.text.strip():
                            rect = get_shape_rect(sub)
                            text_items.append((sub, rect))
            except:
                pass
    
    # 为每个文字框判断应使用什么文字颜色
    for txt_shape, txt_rect in text_items:
        # 先检查：这个文字框自身是否有填充色？
        own_fill = get_fill_hex(txt_shape)
        if own_fill:
            # 自身有填充：根据填充色深浅
            txt_color = C_WHITE if is_dark(own_fill) else C_DARK_TEXT
        else:
            # 自身无填充：找到它"位于"的色块
            txt_color = C_DARK_TEXT  # 默认深灰
            for blk_shape, blk_hex, blk_rect in color_blocks:
                if blk_rect and txt_rect and rect_contains(blk_rect, txt_rect):
                    txt_color = C_WHITE if is_dark(blk_hex) else C_DARK_TEXT
                    break
        
        # 应用文字颜色
        try:
            set_text_rgb(txt_shape.text_frame, txt_color)
        except:
            pass

def main():
    prs = Presentation(PPT_PATH)
    print(f"处理 {len(prs.slides)} 张幻灯片...")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  幻灯片 {i}")
        process_slide(slide)
    prs.save(OUT_PATH)
    print(f"\n✅ 已保存：{OUT_PATH}")

if __name__ == "__main__":
    main()
