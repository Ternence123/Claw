#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""SPC培训PPT – 统一配色美化 + 名词解释页（匹配原模板风格）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil, os, zipfile, lxml.etree as ET_lxml

INPUT  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUTPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_美化版.pptx"

# ── 从原PPT提取的配色 ──────────────────────────
# 标题栏深蓝
TITLE_BG   = RGBColor(0x1E, 0x3A, 0x5F)
# 强调色（青蓝）
ACCENT     = RGBColor(0x08, 0x91, 0xB2)
# 主题色（来自theme1.xml）
DK2        = RGBColor(0x44, 0x54, 0x6A)  # dk2 深蓝灰
LT2        = RGBColor(0xE7, 0xE6, 0xE6)          # lt2 浅灰
A1_BLUE    = RGBColor(0x44, 0x72, 0xC4)          # accent1 蓝
A2_ORANGE  = RGBColor(0xED, 0x7D, 0x31)          # accent2 橙
A5_LBLUE   = RGBColor(0x5B, 0x9B, 0xD5)          # accent5 浅蓝
A6_GREEN   = RGBColor(0x70, 0xAD, 0x47)          # accent6 绿
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1F, 0x29, 0x3A)          # 深色文字
LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF8)          # 卡片浅底

CARD_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
CARD_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
CARD_TEAL  = RGBColor(0xE0, 0xF7, 0xFA)

def clear_placeholder_shapes(slide):
    """删除幻灯片中所有占位符形状"""
    to_remove = []
    for s in slide.shapes:
        if s.is_placeholder:
            to_remove.append(s)
    for s in to_remove:
        s.element.getparent().remove(s.element)

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=0.5):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)

# ══════════════════════════════════════════════════════
#  制作名词解释页（匹配原PPT风格）
# ══════════════════════════════════════════════════════

def make_term_slide(prs, page_no, total_pages, terms):
    """
    制作一页名词解释，风格与原PPT一致：
    - 顶部：深蓝色标题条 + 白色标题文字
    - 内容：卡片式布局，3列 x 3行 = 9个/页
    """
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    clear_placeholder_shapes(slide)

    W = Inches(13.33)
    H = Inches(7.5)

    # ── 顶部标题条（匹配原PPT：深蓝背景） ──────────
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.1))
    set_shape_fill(title_bar, TITLE_BG)
    title_bar.line.fill.background()

    # 标题文字
    title_text = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12), Inches(0.85))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    if total_pages > 1:
        p.text = f"  SPC 核心名词与符号解释（{page_no}/{total_pages}）"
    else:
        p.text = "  SPC 核心名词与符号解释"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ── 左侧窄竖条（匹配原PPT风格） ────────────────
    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), H)
    set_shape_fill(left_bar, ACCENT)
    left_bar.line.fill.background()

    # ── 底部分隔线 ─────────────────────────────────
    footer_line = slide.shapes.add_shape(1, Inches(0.4), Inches(7.35), W - Inches(0.8), Inches(0.01))
    set_shape_fill(footer_line, LT2)
    footer_line.line.fill.background()

    # ── 3列 x 3行 卡片布局 ──────────────────────
    col_x     = [Inches(0.25), Inches(4.55), Inches(8.85)]
    row_y     = [Inches(1.25), Inches(2.85), Inches(4.45)]
    CARD_W    = Inches(4.15)
    CARD_H    = Inches(1.45)
    BG_COLORS = [CARD_GREEN, CARD_BLUE, CARD_TEAL]

    for idx, (sym, name, meaning, formula) in enumerate(terms):
        col = idx % 3
        row = idx // 3
        if row >= 3:
            break   # 最多9个/页

        cx = col_x[col]
        cy = row_y[row]
        bg_color = BG_COLORS[(col + row) % len(BG_COLORS)]

        # 卡片背景
        card = slide.shapes.add_shape(1, cx, cy, CARD_W, CARD_H)
        set_shape_fill(card, bg_color)
        set_shape_line(card, ACCENT, 0.75)

        # 顶部左侧：符号（大号加粗，强调色）
        sym_box = slide.shapes.add_textbox(
            cx + Inches(0.12), cy + Inches(0.07),
            Inches(1.6), Inches(0.32))
        p = sym_box.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_BG

        # 顶部右侧：名称
        name_box = slide.shapes.add_textbox(
            cx + Inches(1.78), cy + Inches(0.09),
            CARD_W - Inches(1.9), Inches(0.30))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(8.5)
        p.font.color.rgb = DARK_TEXT

        # 中部：含义
        mean_box = slide.shapes.add_textbox(
            cx + Inches(0.12), cy + Inches(0.46),
            CARD_W - Inches(0.22), Inches(0.38))
        p = mean_box.text_frame.paragraphs[0]
        p.text = meaning
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_TEXT

        # 底部：公式（斜体，蓝色）
        if formula:
            fbox = slide.shapes.add_textbox(
                cx + Inches(0.12), cy + Inches(0.88),
                CARD_W - Inches(0.22), Inches(0.50))
            # 允许文字换行
            tf2 = fbox.text_frame
            tf2.word_wrap = True
            p = tf2.paragraphs[0]
            p.text = formula
            p.font.size = Pt(7.5)
            p.font.italic = True
            p.font.color.rgb = A1_BLUE

    return slide

# ══════════════════════════════════════════════════════
#  用 XML 操作把追加的页移到指定位置
# ══════════════════════════════════════════════════════

def reorder_slides(input_path, output_path, positions):
    """
    positions: list of (src_index, dst_index)
    将 input_path 的幻灯片重新排列后保存到 output_path
    """
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    # 读取 presentation.xml
    with zipfile.ZipFile(input_path, 'r') as z:
        pres_bytes = z.read('ppt/presentation.xml')
        all_slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide')]

    root = ET_lxml.fromstring(pres_bytes)
    sld_id_lst = root.find('p:sldIdLst', ns)
    sld_ids = sld_id_lst.findall('p:sldId', ns)

    # 要移动的页（从后往前删，避免索引变化）
    elems = []
    for dst, src in sorted(positions, key=lambda x: x[1], reverse=True):
        elem = sld_ids[src]
        elems.append((dst, elem))
        sld_id_lst.remove(elem)

    # 插入到目标位置
    for dst, elem in elems:
        sld_id_lst.insert(dst, elem)

    new_bytes = ET_lxml.tostring(root, xml_declaration=True, encoding='UTF-8')

    tmp = output_path + '.tmp.zip'
    with zipfile.ZipFile(input_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w') as zout:
            for item in zin.namelist():
                if item == 'ppt/presentation.xml':
                    zout.writestr(item, new_bytes)
                else:
                    zout.writestr(item, zin.read(item))
    os.replace(tmp, output_path)
    print(f"  ✅ 页码调整完成")

# ══════════════════════════════════════════════════════
#  main
# ══════════════════════════════════════════════════════

def main():
    prs = Presentation(INPUT)
    orig_n = len(prs.slides)
    print(f"原文件页数: {orig_n}")

    # 17 个术语
    all_terms = [
        ("Σ  (求和)",   "求和符号  Summation",
         "对一系列数值进行累加运算",               "Σxᵢ = x₁ + x₂ + … + xₙ"),
        ("x̄  (X-Bar)", "样本均值  Sample Mean",
         "单个子组内数据的平均值",                 "x̄ = Σxᵢ / n"),
        ("X̄̄  (X-DBar)", "总均值  Grand Mean",
         "所有子组均值的总平均（CL中心线）",       "X̄̄ = (x̄₁+x̄₂+…+x̄k)/k"),
        ("R",          "极差  Range",
         "单个子组内最大值与最小值之差",           "R = xₘₐₓ − xₘᵢₙ"),
        ("R̄  (R-Bar)", "平均极差  Avg Range",
         "所有子组极差 R 的平均值",               "R̄ = (R₁+R₂+…+Rk)/k"),
        ("σ  (Sigma)",  "总体标准差",
         "描述总体数据离散程度的参数",             "σ = √[Σ(xᵢ−μ)²/N]"),
        ("σ̂  (Sigma Hat)", "估计标准差",
         "用样本数据估算的总体标准差",             "σ̂ = R̄/d₂  或  σ̂ = S̄/c₄"),
        ("s / S̄",     "样本标准差 / 平均样本标准差",
         "样本公式算出的标准差及其均值",           "S̄ = (s₁+s₂+…+sk)/k"),
        ("USL",        "规格上限  Upper Spec Limit",
         "客户/设计允许的最大值",                 "由产品设计决定（非统计值）"),
        ("LSL",        "规格下限  Lower Spec Limit",
         "客户/设计允许的最小值",                 "由产品设计决定（非统计值）"),
        ("UCL",        "上控制限  Upper Control Limit",
         "控制图预警上限；超出说明过程异常",       "UCL = X̄̄ + A₂·R̄"),
        ("CL",         "中心线  Center Line",
         "控制图基准线，取总均值 X̄̄",            "CL = X̄̄"),
        ("LCL",        "下控制限  Lower Control Limit",
         "控制图预警下限；低于说明过程异常",       "LCL = X̄̄ − A₂·R̄"),
        ("Cp",         "过程能力指数  Process Capability",
         "衡量过程潜在能力（假设均值无偏移）",     "Cp = (USL−LSL)/(6σ̂)"),
        ("Cpk",        "实际过程能力指数",
         "衡量过程实际能力（考虑均值偏移）",       "Cpk = min[(USL−μ)/3σ̂, (μ−LSL)/3σ̂]"),
        ("Pp / Ppk",   "过程性能指数  Process Performance",
         "用长期总体数据计算",                     "Pp = (USL−LSL)/(6σₜₒₜₐₗ)"),
        ("3σ 原理",     "三西格玛法则  Three-Sigma Rule",
         "正态分布中 99.73% 数据落在 μ±3σ 内",  "P(μ±3σ) = 99.73%"),
    ]

    n_terms = len(all_terms)
    # 每页最多9个，分2页（9 + 8）
    per_page = 9
    n_pages  = (n_terms + per_page - 1) // per_page

    print(f"术语总数: {n_terms}，分 {n_pages} 页")
    added_indices = []
    for p_no in range(n_pages):
        chunk = all_terms[p_no * per_page : (p_no + 1) * per_page]
        print(f"  第{p_no+1}页：{len(chunk)} 个术语")
        slide = make_term_slide(prs, p_no + 1, n_pages, chunk)
        added_indices.append(len(prs.slides) - 1)   # 记录新页的索引

    # 保存到临时文件
    tmp = INPUT + '.tmp.pptx'
    prs.save(tmp)
    print(f"已追加 {n_pages} 页，正在调整页码顺序...")

    # 将追加的页移到第3页位置（索引=2）
    positions = [(2 + i, added_indices[i]) for i in range(len(added_indices))]
    reorder_slides(tmp, OUTPUT, positions)
    os.remove(tmp)

    # 验证
    prs2 = Presentation(OUTPUT)
    print(f"\n✅ 已生成: {OUTPUT}")
    print(f"   总页数: {len(prs2.slides)}")
    for i, s in enumerate(prs2.slides, 1):
        title = ''
        for sh in s.shapes:
            if sh.has_text_frame and sh.text.strip():
                t = sh.text.strip()
                if len(t) > 3:
                    title = t[:35]
                    break
        mark = "  ← 📖 名词解释" if "名词解释" in title else ""
        print(f"   第{i}页: {title}{mark}")

if __name__ == "__main__":
    main()
