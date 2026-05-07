#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""在SPC培训PPT中插入名词解释页，解释所有符号含义和运算方式"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

INPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUTPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v3.pptx"

# 绿色/蓝色主色
GREEN = RGBColor(0x2E, 0x8B, 0x57)   # 深绿
BLUE  = RGBColor(0x1E, 0x3A, 0x8A)   # 深蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
DARK_TEXT   = RGBColor(0x21, 0x21, 0x21)

def add_term_slide(prs):
    """在第3页位置插入名词解释页"""
    
    # 使用空白布局（index=6 通常是空白）
    blank_layout = prs.slide_layouts[6]  # 空白
    # 在索引2（第3页）插入
    slide = prs.slides.add_slide(blank_layout)
    
    # 移动到第3位（索引2）
    # python-pptx 不支持直接插入，只能追加后再排序
    # 改用：在末尾添加后，通过重新构建来插入
    # 更简单：直接加到第3页位置，我们重新生成slides顺序
    
    return slide

def build_explanation_slide(prs):
    """构建名词解释幻灯片，放在第3页位置"""
    
    # 读取所有现有幻灯片
    slides_data = []
    for s in prs.slides:
        slides_data.append(s)
    
    # 新建presentation来重构顺序
    new_prs = Presentation(INPUT)
    
    # 删除所有幻灯片（通过shapes数判断，但python-pptx不能直接删除）
    # 改用：新建一个空白presentation，逐页复制
    
    # 更简单的方法：直接在第3页位置追加一页，然后调整顺序
    # 实际上python-pptx不支持幻灯片重排，我们直接修改第3页后面的内容
    
    # 策略：找到"培训内容"页（第2页），在其后插入一页
    # 由于不能插入，我们在最后添加，然后告知用户
    
    # 改用直接编辑：在现有第2页后新建一页（追加到末尾），然后修改标题
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    
    return new_slide

def make_term_slide(slide):
    """制作名词解释页内容"""
    
    # 背景色块（顶部标题区）
    title_bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = GREEN
    title_bg.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "SPC 核心名词与符号解释"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # 两列布局的术语解释
    terms = [
        # (符号, 名称, 含义, 计算公式)
        ("σ (Sigma)", "标准差/西格玛", 
         "描述数据分布的离散程度，σ越小表示过程越稳定",
         "σ = √[Σ(xi - x̄)²/(n-1)]"),
        
        ("x̄ (X-Bar)", "样本均值",
         "子组内数据的平均值，反映过程中心位置",
         "x̄ = (x₁+x₂+...+xn)/n"),
        
        ("R", "极差",
         "子组内最大值与最小值之差，用于估算过程变异",
         "R = xₘₐₓ - xₘᵢₙ"),
        
        ("USL", "规格上限 (Upper Spec Limit)",
         "客户或设计允许的最大值，超出为不合格",
         "由产品设计决定，非统计数据"),
        
        ("LSL", "规格下限 (Lower Spec Limit)",
         "客户或设计允许的最小值，低于为不合格",
         "由产品设计决定，非统计数据"),
        
        ("UCL", "上控制限 (Upper Control Limit)",
         "控制图预警上限，超出说明过程异常",
         "UCL = x̄̄ + A₂·R̄  (均值图)"),
        
        ("CL", "中心线 (Center Line)",
         "控制图基准线，通常取过程均值",
         "CL = x̄̄ (总均值)"),
        
        ("LCL", "下控制限 (Lower Control Limit)",
         "控制图预警下限，低于说明过程异常",
         "LCL = x̄̄ - A₂·R̄  (均值图)"),
        
        ("Cp", "过程能力指数",
         "衡量过程潜在能力（未考虑偏移）",
         "Cp = (USL-LSL)/(6σ)"),
        
        ("Cpk", "过程能力指数（实际）",
         "衡量过程实际能力（考虑均值偏移）",
         "Cpk = min[(USL-μ)/3σ, (μ-LSL)/3σ]"),
        
        ("Pp / Ppk", "过程性能指数",
         "用长期数据计算的能力指数（总体σ）",
         "Pp = (USL-LSL)/(6σ总体)"),
        
        ("3σ 原理", "三西格玛法则",
         "正态分布中99.73%的数据落在均值±3σ范围内",
         "P(μ±3σ) = 99.73%"),
    ]
    
    # 分两列显示
    col1_x = Inches(0.3)
    col2_x = Inches(6.8)
    start_y = Inches(1.4)
    row_h = Inches(0.72)
    
    for i, (symbol, name, meaning, formula) in enumerate(terms):
        col_x = col1_x if i < 6 else col2_x
        row_idx = i if i < 6 else i - 6
        item_y = start_y + row_idx * row_h
        
        # 背景卡片
        card = slide.shapes.add_shape(
            1,
            col_x, item_y,
            Inches(6.2), row_h - Inches(0.05)
        )
        # 交替颜色
        if row_idx % 2 == 0:
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_GREEN
        else:
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BLUE
        card.line.color.rgb = GREEN
        card.line.width = Pt(0.5)
        
        # 符号（粗体）
        sym_box = slide.shapes.add_textbox(
            col_x + Inches(0.1), item_y + Inches(0.05),
            Inches(1.5), Inches(0.25)
        )
        p = sym_box.text_frame.paragraphs[0]
        p.text = symbol
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        # 名称
        name_box = slide.shapes.add_textbox(
            col_x + Inches(1.6), item_y + Inches(0.05),
            Inches(4.5), Inches(0.25)
        )
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_TEXT
        
        # 含义
        mean_box = slide.shapes.add_textbox(
            col_x + Inches(0.1), item_y + Inches(0.3),
            Inches(5.9), Inches(0.2)
        )
        p = mean_box.text_frame.paragraphs[0]
        p.text = meaning
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_TEXT
        
        # 公式
        if formula:
            form_box = slide.shapes.add_textbox(
                col_x + Inches(0.1), item_y + Inches(0.5),
                Inches(5.9), Inches(0.2)
            )
            p = form_box.text_frame.paragraphs[0]
            p.text = formula
            p.font.size = Pt(8)
            p.font.italic = True
            p.font.color.rgb = BLUE

def main():
    prs = Presentation(INPUT)
    print(f"当前页数: {len(prs.slides)}")
    
    # 构建解释页
    new_slide = build_explanation_slide(prs)
    make_term_slide(new_slide)
    
    # 调整：由于python-pptx不能插入幻灯片，我们需要重建
    # 更简洁的方案：直接修改，把解释页放到第11页位置（要点速记卡之前）
    # 我们先保存，然后告知用户
    
    # 实际上我们需要在第3页后插入，这里用另一个方法：
    # 保存后再用另一个脚本重排，或者我们就放在最后
    
    # 修改方案：找到第11页，把它的内容改成名词解释，原第11页内容后移
    # 但这样会丢失内容
    
    # 最终方案：直接追加到末尾，然后在标题注明"附录"
    
    prs.save(OUTPUT)
    print(f"✅ 已保存: {OUTPUT}")
    print(f"   新增一页「SPC核心名词与符号解释」")
    print(f"   共 {len(prs.slides)} 页（注：新页追加在末尾，建议手动调整到靠前位置）")

if __name__ == "__main__":
    main()
