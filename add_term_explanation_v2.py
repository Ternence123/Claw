#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""在SPC培训PPT第3页位置插入名词解释页"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil
import os
import zipfile
import xml.etree.ElementTree as ET

INPUT  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUTPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_名词解释.pptx"

# 配色
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
BLUE   = RGBColor(0x1E, 0x3A, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
LBLUE  = RGBColor(0xE3, 0xF2, 0xFD)
DARK   = RGBColor(0x21, 0x21, 0x21)

def make_slide_with_terms(prs):
    """制作名词解释页，返回slide对象"""
    layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(layout)
    
    # 顶部绿色标题条
    title_rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.1)
    )
    title_rect.fill.solid()
    title_rect.fill.fore_color.rgb = GREEN
    title_rect.line.fill.background()
    
    # 标题文字
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12), Inches(0.7))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = "📖 SPC 核心名词与符号解释"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 术语数据：(符号, 全称, 含义, 公式)
    terms = [
        ("σ (Sigma)", "标准差 Standard Deviation",
         "数据离散程度的度量，σ越小过程越稳定",
         "σ = √[ Σ(xᵢ−x̄)² / (n−1) ]"),
        
        ("x̄ (X-Bar)", "样本均值 Sample Mean",
         "子组内数据的平均值，反映过程中心位置",
         "x̄ = (x₁+x₂+...+xₙ) / n"),
        
        ("R", "极差 Range",
         "子组内最大值与最小值之差，用于估算过程变异",
         "R = xₘₐₓ − xₘᵢₙ"),
        
        ("USL", "规格上限 Upper Spec Limit",
         "客户/设计允许的最大值，超出即不合格品",
         "由产品设计决定（非统计计算）"),
        
        ("LSL", "规格下限 Lower Spec Limit",
         "客户/设计允许的最小值，低于即不合格品",
         "由产品设计决定（非统计计算）"),
        
        ("UCL", "上控制限 Upper Control Limit",
         "控制图预警上限，超出说明过程出现异常",
         "UCL = x̄̄ + A₂·R̄  （均值图）"),
        
        ("CL", "中心线 Center Line",
         "控制图基准线，通常取所有子组均值的总均值",
         "CL = x̄̄  （总均值）"),
        
        ("LCL", "下控制限 Lower Control Limit",
         "控制图预警下限，低于说明过程出现异常",
         "LCL = x̄̄ − A₂·R̄  （均值图）"),
        
        ("Cp", "过程能力指数 Process Capability",
         "衡量过程潜在能力（假设均值无偏移）",
         "Cp = (USL − LSL) / (6σ̂)"),
        
        ("Cpk", "实际过程能力指数",
         "衡量过程实际能力（考虑均值偏移）",
         "Cpk = min[(USL−μ)/3σ, (μ−LSL)/3σ]"),
        
        ("Pp / Ppk", "过程性能指数 Process Performance",
         "用长期总体数据计算（使用总标准差）",
         "Pp = (USL−LSL) / (6σₜₒₜₐₗ)"),
        
        ("3σ 原理", "三西格玛法则 Three Sigma Rule",
         "正态分布中99.73%的数据落在均值±3σ范围内",
         "P(μ±3σ) = 99.73%"),
    ]
    
    # 两列卡片式布局
    col1_x = Inches(0.2)
    col2_x = Inches(6.8)
    start_y = Inches(1.25)
    row_h = Inches(0.88)
    
    for i, (sym, full_name, meaning, formula) in enumerate(terms):
        col_x = col1_x if i < 6 else col2_x
        row = i if i < 6 else i - 6
        y = start_y + row * row_h
        
        # 卡片背景
        card = slide.shapes.add_shape(
            1,
            col_x, y,
            Inches(6.3), row_h - Inches(0.05)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LGREEN if row % 2 == 0 else LBLUE
        card.line.color.rgb = GREEN
        card.line.width = Pt(0.5)
        
        # 符号（加粗，绿色）
        sym_box = slide.shapes.add_textbox(
            col_x + Inches(0.12), y + Inches(0.06),
            Inches(1.8), Inches(0.28)
        )
        p = sym_box.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        # 全称
        name_box = slide.shapes.add_textbox(
            col_x + Inches(2.0), y + Inches(0.06),
            Inches(4.1), Inches(0.28)
        )
        p = name_box.text_frame.paragraphs[0]
        p.text = full_name
        p.font.size = Pt(9)
        p.font.color.rgb = DARK
        
        # 含义
        mean_box = slide.shapes.add_textbox(
            col_x + Inches(0.12), y + Inches(0.35),
            Inches(6.0), Inches(0.25)
        )
        p = mean_box.text_frame.paragraphs[0]
        p.text = "📌 " + meaning
        p.font.size = Pt(8.5)
        p.font.color.rgb = DARK
        
        # 公式
        if formula:
            fbox = slide.shapes.add_textbox(
                col_x + Inches(0.12), y + Inches(0.6),
                Inches(6.0), Inches(0.25)
            )
            p = fbox.text_frame.paragraphs[0]
            p.text = "🔢 " + formula
            p.font.size = Pt(8)
            p.font.italic = True
            p.font.color.rgb = BLUE
    
    return slide

def reorder_slides(input_path, output_path, insert_at=2):
    """
    用zipfile + xml操作重新排序幻灯片：
    将最后一页（解释页）插入到指定位置
    """
    # 复制文件
    shutil.copy(input_path, output_path)
    
    with zipfile.ZipFile(output_path, 'r') as zin:
        with zipfile.ZipFile(output_path + '.tmp', 'w') as zout:
            for item in zin.namelist():
                content = zin.read(item)
                zout.writestr(item, content)
    
    # 用lxml操作presentation.xml来重排幻灯片顺序
    # 幻灯片ID列表在 ppt/presentation.xml 的 p:sldIdLst 中
    import lxml.etree as ET_lxml
    
    tmp_path = output_path + '.tmp'
    with zipfile.ZipFile(tmp_path, 'r') as z:
        pres_xml = z.read('ppt/presentation.xml')
    
    root = ET_lxml.fromstring(pres_xml)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    
    sld_id_lst = root.find('p:sldIdLst', ns)
    sld_ids = sld_id_lst.findall('p:sldId', ns)
    
    # 最后一页的rId
    last_id = sld_ids[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    last_elem = sld_ids[-1]
    
    # 删除最后一个
    sld_id_lst.remove(last_elem)
    
    # 插入到 insert_at 位置
    sld_id_lst.insert(insert_at, last_elem)
    
    # 写回
    with zipfile.ZipFile(tmp_path, 'r') as zin:
        with zipfile.ZipFile(output_path, 'w') as zout:
            for item in zin.namelist():
                content = zin.read(item)
                if item == 'ppt/presentation.xml':
                    content = ET_lxml.tostring(root, xml_declaration=True, encoding='UTF-8')
                zout.writestr(item, content)
    
    os.remove(tmp_path)

def main():
    # 1. 打开原文件，追加新页（会到最后一页）
    prs = Presentation(INPUT)
    print(f"原文件页数: {len(prs.slides)}")
    
    # 2. 制作解释页（追加到最后）
    new_slide = make_slide_with_terms(prs)
    
    # 3. 先保存，再用XML操作重排幻灯片顺序
    tmp_path = INPUT + '.tmp.pptx'
    prs.save(tmp_path)
    print(f"已追加名词解释页，正在调整页码顺序...")
    
    # 4. 用XML操作将最后一页移到第3页位置
    reorder_slides(tmp_path, OUTPUT, insert_at=2)
    os.remove(tmp_path)
    
    # 验证
    prs2 = Presentation(OUTPUT)
    print(f"\n✅ 已生成: {OUTPUT}")
    print(f"   总页数: {len(prs2.slides)}")
    for i, s in enumerate(prs2.slides, 1):
        title = ''
        for sh in s.shapes:
            if sh.has_text_frame and sh.text.strip():
                title = sh.text.strip()[:35]
                break
        mark = " ← [名词解释]" if "名词解释" in title or "SPC 核心" in title else ""
        print(f"   第{i}页: {title}{mark}")

if __name__ == "__main__":
    main()
