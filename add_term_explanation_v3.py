#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""在SPC培训PPT第3页位置插入名词解释页"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import shutil
import zipfile
import lxml.etree as ET_lxml
import os

INPUT  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUTPUT = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_名词解释.pptx"

GREEN = RGBColor(0x2E, 0x8B, 0x57)
BLUE  = RGBColor(0x1E, 0x3A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREEN= RGBColor(0xE8, 0xF5, 0xE9)
LBLUE = RGBColor(0xE3, 0xF2, 0xFD)
DARK  = RGBColor(0x21, 0x21, 0x21)
GRAY  = RGBColor(0x55, 0x55, 0x55)


def clear_placeholders(slide):
    """清空幻灯片中的所有占位符文字，避免显示默认文字"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                para.text = ""
            # 如果是占位符，隐藏它
            if shape.is_placeholder:
                shape.element.getparent().remove(shape.element)


def make_term_slide(prs):
    """制作名词解释页，追加到末尾"""
    layout = prs.slide_layouts[0]  # 只有 DEFAULT 布局可用
    slide = prs.slides.add_slide(layout)
    clear_placeholders(slide)

    W = Inches(13.33)
    H = Inches(7.5)

    # ── 顶部绿色标题条 ───────────────────────────
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.15))
    rect.fill.solid()
    rect.fill.fore_color.rgb = GREEN
    rect.line.fill.background()

    # 标题文字
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "📖  SPC 核心名词与符号解释"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ── 术语列表 ────────────────────────────────
    terms = [
        ("σ  (Sigma)", "标准差  Standard Deviation",
         "数据离散程度的度量；σ 越小，过程越稳定",
         "σ = √[ Σ(xᵢ − x̄)² / (n−1) ]"),
        ("x̄  (X-Bar)", "样本均值  Sample Mean",
         "子组内数据的平均值，反映过程中心位置",
         "x̄ = (x₁ + x₂ + … + xₙ) / n"),
        ("R", "极差  Range",
         "子组内最大值与最小值之差，用于快速估算变异",
         "R = xₘₐₓ − xₘᵢₙ"),
        ("USL", "规格上限  Upper Spec Limit",
         "客户/设计允许的最大值，超出即不合格品",
         "由产品设计决定（非统计计算值）"),
        ("LSL", "规格下限  Lower Spec Limit",
         "客户/设计允许的最小值，低于即不合格品",
         "由产品设计决定（非统计计算值）"),
        ("UCL", "上控制限  Upper Control Limit",
         "控制图预警上限；超出说明过程出现异常",
         "UCL = x̄̄ + A₂·R̄   （均值图）"),
        ("CL", "中心线  Center Line",
         "控制图基准线，取所有子组均值的总均值",
         "CL = x̄̄  （总均值）"),
        ("LCL", "下控制限  Lower Control Limit",
         "控制图预警下限；低于说明过程出现异常",
         "LCL = x̄̄ − A₂·R̄   （均值图）"),
        ("Cp", "过程能力指数  Process Capability",
         "衡量过程潜在能力（假设过程均值无偏移）",
         "Cp = (USL − LSL) / (6σ̂)"),
        ("Cpk", "实际过程能力指数",
         "衡量过程实际能力（考虑均值偏移）",
         "Cpk = min[(USL−μ)/3σ , (μ−LSL)/3σ]"),
        ("Pp / Ppk", "过程性能指数  Process Performance",
         "用长期总体数据计算（使用总体标准差）",
         "Pp = (USL − LSL) / (6σₜₒₜₐₗ)"),
        ("3σ 原理", "三西格玛法则  Three-Sigma Rule",
         "正态分布中 99.73% 的数据落在均值 ±3σ 范围内",
         "P(μ ± 3σ) = 99.73%"),
    ]

    # ── 两列卡片 ────────────────────────────────
    col1_x = Inches(0.25)
    col2_x = Inches(6.85)
    start_y = Inches(1.3)
    row_h = Inches(0.92)

    for i, (sym, name, meaning, formula) in enumerate(terms):
        cx = col1_x if i < 6 else col2_x
        row = i if i < 6 else i - 6
        cy = start_y + row * row_h

        # 卡片背景
        card = slide.shapes.add_shape(1, cx, cy, Inches(6.3), row_h - Inches(0.06))
        card.fill.solid()
        card.fill.fore_color.rgb = LGREEN if row % 2 == 0 else LBLUE
        card.line.color.rgb = GREEN
        card.line.width = Pt(0.5)

        # 符号（加粗绿色）
        sym_box = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.05),
                                           Inches(1.9), Inches(0.28))
        p = sym_box.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GREEN

        # 名称
        name_box = slide.shapes.add_textbox(cx + Inches(2.1), cy + Inches(0.05),
                                            Inches(4.0), Inches(0.28))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.color.rgb = DARK

        # 含义
        mean_box = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.34),
                                            Inches(6.0), Inches(0.27))
        p = mean_box.text_frame.paragraphs[0]
        p.text = "📌 " + meaning
        p.font.size = Pt(8.5)
        p.font.color.rgb = DARK

        # 公式
        if formula:
            fbox = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.62),
                                            Inches(6.0), Inches(0.25))
            p = fbox.text_frame.paragraphs[0]
            p.text = "🔢 " + formula
            p.font.size = Pt(8)
            p.font.italic = True
            p.font.color.rgb = BLUE

    return slide


def reorder_slides_zip(input_path, output_path, insert_pos=2):
    """
    用 zipfile + lxml 将最后一页移到 insert_pos 位置
    insert_pos: 0-based index，2 表示第3页
    """
    import xml.etree.ElementTree as ET_std

    # 先复制为输出文件，再修改其内部 XML
    shutil.copy(input_path, output_path)

    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    # 读取 presentation.xml
    with zipfile.ZipFile(output_path, 'r') as z:
        pres_bytes = z.read('ppt/presentation.xml')

    root = ET_lxml.fromstring(pres_bytes)
    sld_id_lst = root.find('p:sldIdLst', ns)
    sld_ids = sld_id_lst.findall('p:sldId', ns)

    total = len(sld_ids)
    if insert_pos >= total:
        print(f"⚠ insert_pos={insert_pos} >= 总页数{total}，不调整顺序")
        return

    # 取出最后一页元素
    last_elem = sld_ids[-1]

    # 删除最后一页
    sld_id_lst.remove(last_elem)

    # 插入到指定位置
    sld_id_lst.insert(insert_pos, last_elem)

    # 将修改后的 XML 写回 zip
    new_pres_bytes = ET_lxml.tostring(root, xml_declaration=True, encoding='UTF-8')

    # 重建 zip（修改一个文件）
    import io, os
    tmp_zip = output_path + '.tmp.zip'
    with zipfile.ZipFile(output_path, 'r') as zin:
        with zipfile.ZipFile(tmp_zip, 'w') as zout:
            for item in zin.namelist():
                if item == 'ppt/presentation.xml':
                    zout.writestr(item, new_pres_bytes)
                else:
                    zout.writestr(item, zin.read(item))
    os.replace(tmp_zip, output_path)
    print(f"✅ 页码调整完成：名词解释页已移至第{insert_pos+1}页位置")


def main():
    # 1. 打开原文件
    prs = Presentation(INPUT)
    orig_count = len(prs.slides)
    print(f"原文件页数: {orig_count}")

    # 2. 追加名词解释页（会到最后一页）
    make_term_slide(prs)

    # 3. 保存到临时文件
    tmp_file = INPUT + '.tmp.pptx'
    prs.save(tmp_file)
    print(f"已追加名词解释页，正在调整页码顺序...")

    # 4. 用 zip+XMl 将最后一页移到第3页位置
    reorder_slides_zip(tmp_file, OUTPUT, insert_pos=2)
    os.remove(tmp_file)

    # 5. 验证输出文件
    prs2 = Presentation(OUTPUT)
    print(f"\n✅ 已生成: {OUTPUT}")
    print(f"   总页数: {len(prs2.slides)}")
    for i, s in enumerate(prs2.slides, 1):
        title = ''
        for sh in s.shapes:
            if sh.has_text_frame and sh.text.strip():
                title = sh.text.strip()[:40]
                break
        mark = "  ← 📖 名词解释" if ("名词解释" in title or "SPC 核心" in title) else ""
        print(f"   第{i}页: {title}{mark}")


if __name__ == "__main__":
    main()
