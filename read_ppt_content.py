from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

# 读取现有PPT
prs = Presentation("SPC_3σ原理培训_v2.pptx")

# 提取所有幻灯片的内容
slides_content = []

for i, slide in enumerate(prs.slides):
    slide_data = {
        "slide_num": i + 1,
        "layout": str(slide.slide_layout.name),
        "texts": []
    }
    
    # 提取所有文本框内容
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_data["texts"].append(shape.text.strip())
    
    slides_content.append(slide_data)

# 输出为JSON
print(json.dumps(slides_content, ensure_ascii=False, indent=2))

#  also print raw text for each slide
print("\n" + "="*60)
print("RAW TEXT CONTENT:")
print("="*60)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text.strip())
            print("-" * 40)
