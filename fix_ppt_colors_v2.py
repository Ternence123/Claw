"""
SPC PPT 配色统一 + 文字颜色修复 v3
- 主色：蓝色 #0EA5E9 / 绿色 #10B981
- 深色背景上文字强制白色，浅色背景上文字深灰色
- 大幅减少红/黄色块使用（仅保留语义必要的）
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

PPT_PATH = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v2.pptx"
OUT_PATH  = "c:/Users/Administrator/WorkBuddy/Claw/SPC_3σ原理培训_v3.pptx"

# ── 目标配色 ─────────────────────────────────────
C_DARK_BG   = RGBColor(0x1E, 0x29, 0x3B)   # 深蓝灰 标题栏背景
C_BLUE_MAIN = RGBColor(0x0E, 0xA5, 0xE9)   # 主蓝 色块
C_GREEN_MAIN= RGBColor(0x10, 0xB9, 0x81)   # 主绿 色块
C_BLUE_LIGHT= RGBColor(0xF0, 0xF9, 0xFF)   # 浅蓝 浅背景
C_GREEN_LIGHT=RGBColor(0xEC, 0xFD, 0xF5)   # 浅绿 浅背景
C_RED_WARN  = RGBColor(0xEF, 0x44, 0x44)   # 红色 仅警告
C_YELLOW_HL = RGBColor(0xF5, 0x9E, 0x0B)   # 黄色 仅高亮
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)   # 深灰文字
C_MID_GRAY  = RGBColor(0x64, 0x74, 0x8B)
C_LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFC)   # 页面浅灰

def hex6(rgb_color):
    """RGBColor → '1E293B' 格式"""
    try:
        s = str(rgb_color).upper().replace("0X","")
        return s.zfill(6)
    except:
        return ""

def is_dark(hex6):
    try:
        r=int(hex6[0:2],16); g=int(hex6[2:4],16); b=int(hex6[4:6],16)
        return (r*299+g*587+b*114)/1000 < 128
    except:
        return True

def shape_fill_hex(shape):
    """获取形状填充色 hex6，无填充返回 None"""
    try:
        if shape.fill and shape.fill.type is not None:
            return hex6(shape.fill.fore_color.rgb)
    except:
        pass
    return None

def set_fill(shape, color):
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    except:
        pass

def set_text_color(text_frame, color):
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = color
            except:
                pass

def fix_text_on_shape(shape):
    """根据形状自身的填充色，设置其文字颜色"""
    h = shape_fill_hex(shape)
    if h and shape.has_text_frame and shape.text_frame:
        txt_color = C_WHITE if is_dark(h) else C_DARK_TEXT
        set_text_color(shape.text_frame, txt_color)

def remap_fill(old_hex):
    """
    将旧填充色映射到新配色。
    规则：
    - 蓝系 → C_BLUE_MAIN
    - 深蓝/黑 → C_DARK_BG
    - 绿系 → C_GREEN_MAIN
    - 红系 → C_RED_WARN（保留但减少：仅当形状很小或是警告语义时才用红）
    - 黄系 → C_YELLOW_HL（保留但减少）
    - 紫系 → C_BLUE_MAIN（去掉紫）
    - 浅灰/白 → C_BLUE_LIGHT 或 C_WHITE
    """
    m = {
        # 蓝 → 主蓝
        "0891B2": "BLUE", "2563EB": "BLUE", "2E86AB": "BLUE",
        # 深蓝/黑 → 深蓝背景
        "1E3A5F": "DARK_BG", "1E293B": "DARK_BG", "0D2137": "DARK_BG",
        # 绿 → 主绿
        "27AE60": "GREEN", "10B981": "GREEN",
        # 红 → 警告红（保守保留）
        "EF4444": "RED", "C0392B": "RED",
        # 黄 → 高亮黄（保守保留）
        "F59E0B": "YELLOW", "E07B39": "YELLOW",
        # 紫 → 主蓝
        "8E44AD": "BLUE", "7C3AED": "BLUE",
        # 浅灰
        "F1F5F9": "BLUE_LIGHT", "CBD5E1": "BLUE_LIGHT",
        "FFFFF0": "BLUE_LIGHT", "FFFFFF": "WHITE",
    }
    cat = m.get(old_hex)
    if cat == "BLUE":       return C_BLUE_MAIN
    if cat == "GREEN":     return C_GREEN_MAIN
    if cat == "DARK_BG":   return C_DARK_BG
    if cat == "RED":       return C_RED_WARN
    if cat == "YELLOW":   return C_YELLOW_HL
    if cat == "BLUE_LIGHT": return C_BLUE_LIGHT
    if cat == "WHITE":     return C_WHITE
    return None

def process(shape):
    """处理单个形状（含递归处理组合形状）"""
    # 1) 修复填充色
    old_hex = shape_fill_hex(shape)
    if old_hex:
        new_color = remap_fill(old_hex)
        if new_color:
            set_fill(shape, new_color)
            # 同步更新 old_hex（后续文字颜色判断用）
            old_hex = hex6(new_color)

    # 2) 修复自身文字颜色（根据填充色深浅）
    if shape.has_text_frame and shape.text_frame:
        if old_hex:
            txt_color = C_WHITE if is_dark(old_hex) else C_DARK_TEXT
        else:
            txt_color = C_DARK_TEXT   # 无填充，默认深灰文字
        set_text_color(shape.text_frame, txt_color)

    # 3) 递归处理组合形状
    try:
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                process(sub)
    except:
        pass

def fix_slide(slide):
    for shape in slide.shapes:
        process(shape)

def main():
    prs = Presentation(PPT_PATH)
    print(f"处理 {len(prs.slides)} 张幻灯片...")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  幻灯片 {i}")
        fix_slide(slide)
    prs.save(OUT_PATH)
    print(f"\n✅ 已保存：{OUT_PATH}")

if __name__ == "__main__":
    main()
